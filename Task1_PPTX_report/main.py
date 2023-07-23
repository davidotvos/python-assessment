#!/usr/bin/env python3
import json
import argparse
import os
from pptx.util import Inches
from pptx import Presentation
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO

def validate_json_file(file_path):
    if not file_path.lower().endswith('.json') or not os.path.isfile(file_path):
        raise argparse.ArgumentTypeError('Invalid JSON file: {}'.format(file_path))
    return file_path


def generate_title_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[0]  # Title slide layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def generate_text_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[1]  # Text slide layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def generate_list_slide(presentation, title, items):
    slide_layout = presentation.slide_layouts[1]  # List slide layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    content_frame = slide.placeholders[1].text_frame
    for item in items:
        p = content_frame.add_paragraph()
        p.text = item["text"]
        p.level = item["level"]


def generate_picture_slide(presentation, title, picture_path):

    slide_layout = presentation.slide_layouts[5]  # Picture slide layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Get the slide dimensions
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Get the picture dimensions
    picture = slide.shapes.add_picture(picture_path, 0, 0)
    picture_width = picture.width
    picture_height = picture.height

    # Calculate the position to center the picture
    left = (slide_width - picture_width) // 2
    top = (slide_height - picture_height) // 2

    # Set the picture position
    picture.left = left
    picture.top = top


def read_data_from_dat_file(file_path):
    data = np.loadtxt(file_path, delimiter=';')
    return data[:, 0], data[:, 1]  # Extract first and second columns


def generate_plot_slide(presentation, title, x_data, y_data, configuration):

    # Create the plot using matplotlib
    plt.figure(figsize=(6, 4.5))
    plt.plot(x_data, y_data, linestyle='-')  # Connect dots with a line
    plt.xlabel(configuration["x-label"])
    plt.ylabel(configuration["y-label"])

    # Save the plot as an image in memory
    plot_image = BytesIO()
    plt.savefig(plot_image, format='png')
    plt.close()

    generate_picture_slide(presentation, title, plot_image)


def generate_presentation(config_file):
    with open(config_file) as f:
        config = json.load(f)

    presentation = Presentation()

    for slide_config in config["presentation"]:
        slide_type = slide_config["type"]
        title = slide_config["title"]
        content = slide_config.get("content")
        configuration = slide_config.get("configuration", {})
        items = slide_config.get("content", [])

        if slide_type == "title":
            generate_title_slide(presentation, title, content)
        elif slide_type == "text":
            generate_text_slide(presentation, title, content)
        elif slide_type == "list":
            generate_list_slide(presentation, title, items)
        elif slide_type == "picture":
            generate_picture_slide(presentation, title, content)
        elif slide_type == "plot":
            x_data, y_data = read_data_from_dat_file('sample.dat')
            generate_plot_slide(presentation, title, x_data, y_data ,configuration)
        else:
            raise ValueError(f"Invalid slide type: {slide_type}")

    presentation.save("output.pptx")
    print("Presentation generated successfully.")

def main():

    # Instantiate the parser
    parser = argparse.ArgumentParser(description='Generate a pptx report from a configuration file.')
    parser.add_argument('config_file',type=validate_json_file, help='Path to the JSON configuration file')

    args = parser.parse_args()
    config_file = args.config_file

    generate_presentation(config_file=config_file)
##############################################################################

if __name__ == "__main__":
    main()