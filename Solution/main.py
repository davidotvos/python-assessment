#!/usr/bin/env python3
import json
import argparse
import os
from pptx.util import Inches
from pptx import Presentation
import numpy as np
import matplotlib.pyplot as plt
from io import BytesIO
import logging


# Set up logger configuration
logging.basicConfig(filename='pptx_generator.log', level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s', datefmt='%Y-%m-%d %H:%M:%S')


def validate_json_file(file_path):
    """
    Validate the JSON file path and check if it has a '.json' extension.

    Parameters:
        file_path (str): The path to the JSON configuration file.

    Returns:
        str: The valid file path if it meets the conditions.

    Raises:
        argparse.ArgumentTypeError: If the file path is not valid or does not have a '.json' extension.
    """
    try:
        if not file_path.lower().endswith('.json') or not os.path.isfile(file_path):
            logging.error(f"Invalid JSON file")
            raise argparse.ArgumentTypeError(f'Invalid JSON file: {file_path}')
        
        return file_path
    except Exception as e:
        logging.error(f"Error in validate_json_file: {str(e)}")
        raise e


def generate_title_slide(presentation, title, content):
    """
    Generate a title slide with the given title and content.

    Parameters:
        presentation (pptx.Presentation): The PowerPoint presentation object.
        title (str): The title text for the slide.
        content (str): The content/subtitle text for the slide.
    """
    try:
        slide_layout = presentation.slide_layouts[0]  # Title slide layout
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        logging.info(f"Title slide added")
    except Exception as e:
        logging.error(f"Error in generate_title_slide: {str(e)}")
        raise e


def generate_text_slide(presentation, title, content):
    """
    Generate a text slide with the given title and content.

    Parameters:
        presentation (pptx.Presentation): The PowerPoint presentation object.
        title (str): The title text for the slide.
        content (str): The content/text for the slide.
    """
    try:
        slide_layout = presentation.slide_layouts[1]  # Text slide layout
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        logging.info(f"Text slide added")
    except Exception as e:
        logging.error(f"Error in generate_text_slide: {str(e)}")
        raise e


def generate_list_slide(presentation, title, items):
    """
    Generate a list slide with the given title and items.

    Parameters:
        presentation (pptx.Presentation): The PowerPoint presentation object.
        title (str): The title text for the slide.
        items (list): A list of dictionaries containing 'level' and 'text' keys for list items.
            Example: [{'level': 1, 'text': 'Item 1'}, {'level': 2, 'text': 'Item 1.1'}]
    """
    try:
        slide_layout = presentation.slide_layouts[1]  # List slide layout
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        content_frame = slide.placeholders[1].text_frame
        for item in items:
            if not isinstance(item, dict) or 'level' not in item or 'text' not in item:
                raise ValueError("Invalid list item format. Each item should be a dictionary with 'level' and 'text' keys.")
            p = content_frame.add_paragraph()
            p.text = item["text"]
            p.level = item["level"]

        logging.info(f"List slide added")
    except (ValueError, IndexError, KeyError) as e:
        logging.error(f"Error in generate_list_slide: {str(e)}")
        raise e


def generate_picture_slide(presentation, title, picture_path, log_picture):
    """
    Generate a picture slide with the given title and picture.

    Parameters:
        presentation (pptx.Presentation): The PowerPoint presentation object.
        title (str): The title text for the slide.
        picture_path (str): The file path to the picture to be added to the slide.
    """
    try:
        slide_layout = presentation.slide_layouts[5]  # Picture slide layout
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        # Get the slide dimensions
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Get the picture dimensions
        picture = slide.shapes.add_picture(picture_path, 0, 0)
        picture_width = picture.width
        picture_height = picture.height

        # Calculate the position to center the picture
        left = (slide_width - picture_width) // 2
        top = (slide_height - picture_height) // 2

        # Set the picture position
        picture.left = left
        picture.top = top

        if log_picture:
            logging.info(f"Picture slide added")

    except (OSError, ValueError, Exception) as e:
        logging.error(f"Error in generate_picture_slide: {str(e)}")
        raise e
    

def read_data_from_dat_file(file_path):
    """
    Read data from a .dat file and extract the first and second columns.

    Parameters:
        file_path (str): The path to the .dat file.

    Returns:
        tuple: A tuple containing the x-axis data and y-axis data as numpy arrays.
    """
    try:
        data = np.loadtxt(file_path, delimiter=';')
        x_data = data[:, 0]
        y_data = data[:, 1]
        return x_data, y_data
    except (OSError, ValueError, Exception) as e:
        logging.error(f"Error in read_data_from_dat_file: {str(e)}")
        raise e


def generate_plot_slide(presentation, title, x_data, y_data, configuration):
    """
    Generate a plot slide with the given title and plot data.

    Parameters:
        presentation (pptx.Presentation): The PowerPoint presentation object.
        title (str): The title text for the slide.
        x_data (array-like): The x-axis data for the plot.
        y_data (array-like): The y-axis data for the plot.
        configuration (dict): A dictionary containing plot configuration data, such as 'x-label' and 'y-label'.
    """
    try:
        # Create the plot using matplotlib
        plt.figure(figsize=(6, 4.5))
        plt.plot(x_data, y_data, linestyle='-')  # Connect dots with a line
        plt.xlabel(configuration["x-label"])
        plt.ylabel(configuration["y-label"])

        # Save the plot as an image in memory
        plot_image = BytesIO()
        plt.savefig(plot_image, format='png')
        plt.close()

        generate_picture_slide(presentation, title, plot_image, log_picture=False)
        logging.info(f"Plot slide added")

    except Exception as e:
        logging.error(f"Error in generate_plot_slide: {str(e)}")
        raise e


def generate_presentation(config_file):
    """
    Generate a PowerPoint presentation based on the configuration file.

    Parameters:
        config_file (str): The path to the JSON configuration file.
    """
    try:
        with open(config_file) as f:
            config = json.load(f)

        presentation = Presentation()

        for slide_config in config["presentation"]:
            slide_type = slide_config["type"]
            title = slide_config["title"]
            content = slide_config.get("content")
            configuration = slide_config.get("configuration", {})
            items = slide_config.get("content", [])

            if slide_type == "title":
                generate_title_slide(presentation, title, content)
            elif slide_type == "text":
                generate_text_slide(presentation, title, content)
            elif slide_type == "list":
                generate_list_slide(presentation, title, items)
            elif slide_type == "picture":
                generate_picture_slide(presentation, title, content, log_picture=True)
            elif slide_type == "plot":
                x_data, y_data = read_data_from_dat_file('sample.dat')
                generate_plot_slide(presentation, title, x_data, y_data, configuration)
            else:
                raise ValueError(f"Invalid slide type: {slide_type}")

        presentation.save("output.pptx")
        logging.info("Presentation generated successfully and saved as output.pptx")
    except (FileNotFoundError, json.JSONDecodeError, KeyError, ValueError, Exception) as e:
        logging.error(f"Error in generate_presentation: {str(e)}")
        raise e


def main():

    # Instantiate the parser
    parser = argparse.ArgumentParser(description='Generate a pptx report from a configuration file.')
    parser.add_argument('config_file',type=validate_json_file, help='Path to the JSON configuration file')

    args = parser.parse_args()
    config_file = args.config_file

    generate_presentation(config_file=config_file)
##############################################################################

if __name__ == "__main__":
    main()